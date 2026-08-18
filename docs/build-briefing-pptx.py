"""Builds the briefing as a native PowerPoint deck.

Text, tables and diagram shapes are real PowerPoint objects rather than pictures,
so the deck can be edited in the meeting. The palette and the type hierarchy follow
the HTML version; anything changed in one has to be changed in the other.
"""
from pptx import Presentation
from pptx.util import Emu, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml import parse_xml
from pptx.enum.dml import MSO_LINE_DASH_STYLE

INK   = RGBColor(0x12, 0x21, 0x1F)
MUTED = RGBColor(0x5A, 0x6C, 0x6A)
LINE  = RGBColor(0xD3, 0xDA, 0xD6)
BRAND = RGBColor(0x0D, 0x4F, 0x4A)
WASH  = RGBColor(0xE3, 0xED, 0xEA)
SURF2 = RGBColor(0xEA, 0xEE, 0xEB)
WARN  = RGBColor(0xB0, 0x65, 0x1A)
WARNW = RGBColor(0xF7, 0xEE, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SANS = "맑은 고딕"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Cm(33.867)   # 16:9
prs.slide_height = Cm(19.05)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height
M = Cm(1.6)


def text(slide, x, y, w, h, runs, size=12, color=INK, bold=False, font=SANS,
         align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = runs if isinstance(runs, list) else [(runs, {})]
    for i, (content, opts) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("spacing", spacing)
        if i: p.space_before = Pt(opts.get("before", 4))
        r = p.add_run(); r.text = content
        f = r.font
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.name = opts.get("font", font)
        f.color.rgb = opts.get("color", color)
    return box


def box(slide, x, y, w, h, fill=SURF2, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        int(x), int(y), int(w), int(h))
    shape.adjustments[0] = 0.04 if radius else 0
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    shape.text_frame.word_wrap = True
    return shape


def card(slide, x, y, w, h, title, body, fill=SURF2, title_color=BRAND):
    box(slide, x, y, w, h, fill=fill)
    text(slide, x + Cm(0.45), y + Cm(0.35), w - Cm(0.9), Cm(0.6), title, size=11.5,
         bold=True, color=title_color)
    text(slide, x + Cm(0.45), y + Cm(1.0), w - Cm(0.9), h - Cm(1.3), body, size=10,
         color=MUTED, spacing=1.3)


def rail(slide, label, number):
    text(slide, M, Cm(0.9), Cm(16), Cm(0.6), label.upper(), size=9, bold=True,
         color=BRAND, font=MONO)
    text(slide, W - M - Cm(6), Cm(0.9), Cm(6), Cm(0.6), number, size=9, color=MUTED,
         font=MONO, align=PP_ALIGN.RIGHT)
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(M), Cm(1.55), int(W - M), Cm(1.55))
    ln.line.color.rgb = LINE; ln.line.width = Pt(0.75)


def heading(slide, s, y=Cm(2.1), size=24):
    text(slide, M, y, W - 2 * M, Cm(1.4), s, size=size, bold=True, color=INK, spacing=1.1)


def table(slide, x, y, w, rows, col_w, header=True, size=9.5):
    n, cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n, cols, int(x), int(y), int(w), int(Cm(0.62) * n))
    tbl = shape.table
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Emu(int(w * cw))
    for r, row in enumerate(rows):
        tbl.rows[r].height = Cm(0.62)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Cm(0.2)
            cell.margin_top = cell.margin_bottom = Cm(0.08)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = str(val)
            run.font.size = Pt(size - (1 if header and r == 0 else 0))
            run.font.name = MONO if (isinstance(val, str) and val.startswith(("(", "GET", "1.", "0.", "TAB", "X-", "MED", "ISO", "C-"))) else SANS
            run.font.bold = header and r == 0
            run.font.color.rgb = MUTED if header and r == 0 else INK
    return shape


def arrow(slide, x1, y1, x2, y2, color=INK, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    c.line.color.rgb = color; c.line.width = Pt(1.25)
    if dashed:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    ln = c.line._get_or_add_ln()
    ln.append(parse_xml(
        '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'type="triangle" w="med" len="med"/>'))
    return c


# ---------------------------------------------------------------- slide 1
s = prs.slides.add_slide(BLANK)
rail(s, "MedMission", "01 / 05")
text(s, M, Cm(2.2), W - 2 * M, Cm(0.6), "결핵 검진 현장 · 태블릿 문진에서 X-RAY 워크리스트까지",
     size=9.5, color=MUTED, font=MONO)
heading(s, "태블릿과 브릿지, 두 프로그램 개요", y=Cm(3.0), size=30)
text(s, M, Cm(4.9), Cm(20), Cm(1.4),
     "태블릿 문진이 랩톱을 거쳐 콘솔 워크리스트로 그대로 올라감.\n촬영 담당이 환자 정보를 다시 입력하지 않는 것이 목표.",
     size=13, color=MUTED)

y0 = Cm(7.2)
text(s, M, y0 - Cm(0.7), Cm(6), Cm(0.5), "현재", size=9.5, color=MUTED, font=MONO)
box(s, M, y0, Cm(5.2), Cm(1.5), fill=WHITE, line=LINE)
text(s, M, y0 + Cm(0.45), Cm(5.2), Cm(0.6), "종이 문진", size=11, align=PP_ALIGN.CENTER)
box(s, M + Cm(7.6), y0, Cm(6.6), Cm(1.5), fill=WARNW, line=WARN)
text(s, M + Cm(7.6), y0 + Cm(0.25), Cm(6.6), Cm(1.1),
     [("촬영실에서 재입력", {"size": 11, "bold": True, "color": WARN, "align": PP_ALIGN.CENTER}),
      ("사람 손 · 오타 · 대기", {"size": 9, "color": WARN, "align": PP_ALIGN.CENTER})])
box(s, M + Cm(16.4), y0, Cm(5.2), Cm(1.5), fill=WHITE, line=LINE)
text(s, M + Cm(16.4), y0 + Cm(0.45), Cm(5.2), Cm(0.6), "촬영", size=11, align=PP_ALIGN.CENTER)
arrow(s, M + Cm(5.4), y0 + Cm(0.75), M + Cm(7.4), y0 + Cm(0.75))
arrow(s, M + Cm(14.4), y0 + Cm(0.75), M + Cm(16.2), y0 + Cm(0.75), color=WARN)
text(s, M + Cm(22.2), y0 + Cm(0.3), Cm(7), Cm(1.1),
     "문진 내용은 판독까지 가지 못함", size=9.5, color=MUTED)

y1 = Cm(10.0)
text(s, M, y1 - Cm(0.7), Cm(6), Cm(0.5), "변경 후", size=9.5, color=BRAND, font=MONO)
box(s, M, y1, Cm(5.2), Cm(1.5), fill=WASH, line=LINE)
text(s, M, y1 + Cm(0.45), Cm(5.2), Cm(0.6), "태블릿 문진", size=11, align=PP_ALIGN.CENTER)
box(s, M + Cm(7.6), y1, Cm(6.6), Cm(1.5), fill=WASH, line=LINE)
text(s, M + Cm(7.6), y1 + Cm(0.25), Cm(6.6), Cm(1.1),
     [("랩톱 (브릿지)", {"size": 11, "bold": True, "align": PP_ALIGN.CENTER}),
      ("워크리스트 자동 등재", {"size": 9, "color": MUTED, "align": PP_ALIGN.CENTER})])
box(s, M + Cm(16.4), y1, Cm(5.2), Cm(1.5), fill=WHITE, line=LINE)
text(s, M + Cm(16.4), y1 + Cm(0.45), Cm(5.2), Cm(0.6), "촬영", size=11, align=PP_ALIGN.CENTER)
arrow(s, M + Cm(5.4), y1 + Cm(0.75), M + Cm(7.4), y1 + Cm(0.75))
arrow(s, M + Cm(14.4), y1 + Cm(0.75), M + Cm(16.2), y1 + Cm(0.75))
text(s, M + Cm(22.2), y1 + Cm(0.3), Cm(7), Cm(1.1),
     "문진 답변은 API로 조회 가능", size=9.5, color=BRAND)

text(s, M, Cm(12.6), W - 2 * M, Cm(0.6),
     "없애려는 것은 가운데 한 칸 — 촬영실에서 환자 정보를 사람이 다시 입력하는 단계.",
     size=10, color=MUTED)

cw = int((W - 2 * M - Cm(1.0)) / 3)
for i, (t, b) in enumerate([
        ("전제", "인터넷 없는 폐쇄망. 브릿지는 서버 불필요 — 랩톱 1대 + 같은 Wi-Fi의 태블릿 몇 대. PACS 서버 PC와 달리 네트워크가 흔들려도 태블릿이 문진을 물고 기다림."),
        ("이 회의의 목적", "랩톱 SW가 워크리스트를 조회하고, 촬영 영상의 환자 ID로 문진을 가져가는 부분 합의. 매뉴얼과 실사 160건 분석 완료 — 상대 수정 없이 붙는 상태."),
        ("브릿지가 하지 않는 것", "먼저 호출하지 않음. 영상 다루지 않음. 저장하지 않음. 조회당하는 쪽.")]):
    card(s, M + i * (cw + Cm(0.5)), Cm(13.6), cw, Cm(3.6), t, b,
         fill=WASH if i == 1 else SURF2)

# ---------------------------------------------------------------- slide 2
s = prs.slides.add_slide(BLANK)
rail(s, "전체 구성", "02 / 05")
heading(s, "랩톱 한 대가 가운데")

bx, by = M, Cm(4.6)
box(s, bx, by + Cm(1.0), Cm(5.4), Cm(2.4), fill=WHITE, line=LINE)
text(s, bx, by + Cm(1.5), Cm(5.4), Cm(1.4),
     [("태블릿 ×N", {"size": 12, "bold": True, "align": PP_ALIGN.CENTER}),
      ("Android 문진 앱", {"size": 9.5, "color": MUTED, "align": PP_ALIGN.CENTER})])

mx = bx + Cm(8.4)
box(s, mx, by, Cm(7.4), Cm(4.6), fill=WASH, line=LINE)
text(s, mx, by + Cm(0.4), Cm(7.4), Cm(3.8),
     [("브릿지 (랩톱)", {"size": 12.5, "bold": True, "align": PP_ALIGN.CENTER}),
      ("문진 수신 API · SQLite", {"size": 9.5, "color": MUTED, "align": PP_ALIGN.CENTER, "before": 8}),
      ("DICOM MWL SCP", {"size": 9.5, "color": MUTED, "align": PP_ALIGN.CENTER}),
      ("MPPS SCP", {"size": 9.5, "color": MUTED, "align": PP_ALIGN.CENTER}),
      ("관리 화면 · mDNS 광고", {"size": 9.5, "color": MUTED, "align": PP_ALIGN.CENTER})])

rx = mx + Cm(10.4)
box(s, rx, by, Cm(6.4), Cm(4.6), fill=WHITE, line=WARN)
text(s, rx, by + Cm(0.4), Cm(6.4), Cm(3.8),
     [("MDVizio-X AI", {"size": 12.5, "bold": True, "color": WARN, "align": PP_ALIGN.CENTER}),
      ("MWL SCU · 촬영 · AI 판독", {"size": 9.5, "color": WARN, "align": PP_ALIGN.CENTER, "before": 8}),
      ("MPPS SCU · PACS SCU", {"size": 9.5, "color": WARN, "align": PP_ALIGN.CENTER}),
      ("콘솔과 판독 SW가 같은 프로그램", {"size": 9, "color": MUTED, "align": PP_ALIGN.CENTER, "before": 8})])

arrow(s, bx + Cm(5.5), by + Cm(2.2), mx - Cm(0.1), by + Cm(2.2))
text(s, bx + Cm(5.4), by + Cm(1.5), Cm(3.0), Cm(0.6), "문진 전송", size=9, color=MUTED,
     align=PP_ALIGN.CENTER)
arrow(s, rx - Cm(0.1), by + Cm(1.4), mx + Cm(7.5), by + Cm(1.4), color=WARN)
text(s, mx + Cm(7.6), by + Cm(0.7), Cm(2.6), Cm(0.6), "C-FIND", size=9, color=WARN,
     align=PP_ALIGN.CENTER)
arrow(s, mx + Cm(7.5), by + Cm(3.2), rx - Cm(0.1), by + Cm(3.2), color=WARN, dashed=True)
text(s, mx + Cm(7.6), by + Cm(3.4), Cm(2.6), Cm(0.6), "MPPS", size=9, color=WARN,
     align=PP_ALIGN.CENTER)

text(s, M, Cm(9.8), W - 2 * M, Cm(1.2),
     "브릿지는 먼저 연결하지 않음. 태블릿이 밀어 넣고, MDVizio-X가 가져감. 점선은 MDVizio-X가 보내오는 방향.",
     size=10, color=MUTED)

cw2 = int((W - 2 * M - Cm(0.6)) / 2)
card(s, M, Cm(11.3), cw2, Cm(3.4), "태블릿 앱",
     "7개 섹션 문진, 입력 중 자동 저장. 모든 기록을 태블릿 DB에 먼저 저장하므로 Wi-Fi가 끊겨도 접수 계속 가능. 브릿지는 mDNS로 자동 탐색, 전송 실패 건은 15분 주기 재시도.")
card(s, M + cw2 + Cm(0.6), Cm(11.3), cw2, Cm(3.4), "브릿지",
     "폴더 복사 후 실행 파일 하나. .NET 설치·DB 서버 불필요. MDVizio-X에는 워크리스트로, HTTP를 부를 수 있는 쪽에는 REST로 내어줌. 관리 화면은 그 랩톱 본체에서만 열림.")
text(s, M, Cm(15.2), W - 2 * M, Cm(0.6),
     "Android 8.0+     서명 릴리스 APK     .NET 9 단일 exe     SQLite",
     size=9.5, color=BRAND, font=MONO)

# ---------------------------------------------------------------- slide 3
s = prs.slides.add_slide(BLANK)
rail(s, "흐름과 상태", "03 / 05")
heading(s, "언제 워크리스트에 오르고, 언제 내려가는가")

steps = [("1 · 문진 작성", "태블릿 · 번호 자동 부여", INK),
         ("2 · 랩톱 전송", "실패 시 자동 재시도", INK),
         ("3 · 워크리스트 등재", "콘솔이 C-FIND 조회", WARN),
         ("4 · 촬영", "영상의 환자 ID로 문진 조회", WARN),
         ("5 · 완료 처리", "MPPS면 자동, 아니면 수동", INK)]
sw = int((W - 2 * M) / 5)
ty = Cm(5.4)
for i, (t, sub, col) in enumerate(steps):
    cx = M + i * sw
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx + sw / 2 - Cm(0.18)), int(ty), Cm(0.36), Cm(0.36))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
    dot.shadow.inherit = False
    text(s, cx, ty - Cm(1.0), sw, Cm(0.6), t, size=10.5, bold=True, color=col,
         align=PP_ALIGN.CENTER)
    text(s, cx, ty + Cm(0.6), sw, Cm(1.0), sub, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, cx + sw / 2 + Cm(0.3), ty + Cm(0.18), cx + sw + sw / 2 - Cm(0.3), ty + Cm(0.18),
              color=MUTED)

band = box(s, M + 2 * sw + sw / 2 - Cm(0.5), Cm(8.6), sw * 2 + Cm(1.0), Cm(1.9), fill=WASH, line=LINE)
text(s, M + 2 * sw + sw / 2 - Cm(0.5), Cm(8.9), sw * 2 + Cm(1.0), Cm(1.4),
     [("워크리스트에 노출되는 구간", {"size": 10.5, "bold": True, "align": PP_ALIGN.CENTER}),
      ("Received · InProgress", {"size": 9.5, "color": MUTED, "align": PP_ALIGN.CENTER, "font": MONO})])
box(s, M, Cm(8.6), sw * 2, Cm(1.9), fill=WHITE, line=LINE)
text(s, M, Cm(8.9), sw * 2, Cm(1.4),
     [("아직 랩톱에 없음", {"size": 10.5, "color": MUTED, "align": PP_ALIGN.CENTER}),
      ("태블릿에만 보관", {"size": 9, "color": MUTED, "align": PP_ALIGN.CENTER})])

text(s, M, Cm(11.0), W - 2 * M, Cm(1.2),
     [("5번 이후 상태 — Completed · Cancelled. 워크리스트에서 제외되나 목록에는 잔존", {"size": 10, "color": MUTED}),
      ("실사 160건은 전부 IN PROGRESS로 남아 있었음 — 완료 처리가 한 번도 되지 않았다는 뜻", {"size": 10, "color": WARN})])

card(s, M, Cm(13.2), cw2, Cm(3.2), "150명 규모에서 생기는 문제",
     "완료 처리가 밀리면 목록이 하루 종일 누적. 눈앞 환자를 긴 목록에서 찾아야 함.",
     fill=WARNW, title_color=WARN)
card(s, M + cw2 + Cm(0.6), Cm(13.2), cw2, Cm(3.2), "해법 — MPPS · 구현 완료",
     "MDVizio-X가 촬영 시작·완료를 브릿지에 알려주면 상태가 자동 변경. 설정에서 MPPS 주소만 지정하면 됨 — 워크리스트와 같은 포트·AE.",
     fill=WASH)

# ---------------------------------------------------------------- slide 4
s = prs.slides.add_slide(BLANK)
rail(s, "연동 지점", "04 / 05")
heading(s, "표준 C-FIND로 목록, 문진 전량은 DICOM에 동봉")

colw = int((W - 2 * M - Cm(0.8)) / 2)
text(s, M, Cm(4.2), colw, Cm(0.5), "A · DICOM 워크리스트", size=9, color=MUTED, font=MONO)
table(s, M, Cm(4.8), colw, [
    ("항목", "값"),
    ("AE Title", "MEDMISSION · 검사 안 함"),
    ("포트", "11112 · 랩톱별 상이 · MPPS 동일"),
    ("지원", "C-ECHO, MWL C-FIND, MPPS"),
    ("Modality", "DX · 실사 160건 근거"),
    ("문자셋", "ISO_IR 192 · 상대는 IR 149"),
    ("조회 키", "Patient ID, Name, Accession, Date"),
], col_w=[0.32, 0.68])
text(s, M, Cm(9.7), colw, Cm(2.4),
     "표에 없는 조회 키는 거부하지 않고 무시. 이름·접수번호는 와일드카드와 대소문자 무시 — 접수번호를 손으로 치는 동선이라 관대하게 처리. Windows가 11112를 예약한 랩톱은 12112 등으로 바뀌므로 포트는 현장 확인 필요.",
     size=9.5, color=MUTED)

x2 = M + colw + Cm(0.8)
text(s, x2, Cm(4.2), colw, Cm(0.5), "B · 문진 전량 · 워크리스트에 동봉", size=9, color=MUTED, font=MONO)
table(s, x2, Cm(4.8), colw, [
    ("층", "위치"),
    ("표준 태그", "전화 · 키 · 몸무게 · 흡연"),
    ("판독용 텍스트", "(0010,21B0) 전 항목"),
    ("JSON 원문", "(1001,xx31) UT"),
    ("항목별 name/value", "(1001,xx40) SQ"),
    ("식별자", "recordId = Patient ID"),
], col_w=[0.42, 0.58])
text(s, x2, Cm(9.1), colw, Cm(3.0),
     "MDVizio-X에 HTTP 클라이언트가 없어 REST를 못 씀. 그래서 문진 전 항목을 DICOM에 실어 보냄. 사설 블록은 상대가 AI 결과를 쓰는 (1001) MDAI_PRIVATE_CREATOR 그대로 — 요소 30번대. 상대 수정 0.",
     size=9.5, color=MUTED)

card(s, M, Cm(13.0), colw, Cm(3.4), "매칭은 Patient ID로",
     "recordId = Patient ID (0010,0020) · UUID · 항상 존재하고 유일. 영상에서 이 값을 읽어 문진을 조회하는 것이 권장 경로.",
     fill=WASH)
card(s, x2, Cm(13.0), colw, Cm(3.4), "매칭 기준으로 부적합한 것",
     "no = Accession Number는 비거나 겹칠 수 있음. 환자 이름도 비어 있을 수 있음. 전 구간 평문 HTTP — 폐쇄망 전제로 TLS 보류 중.",
     fill=WARNW, title_color=WARN)

# ---------------------------------------------------------------- slide 5
s = prs.slides.add_slide(BLANK)
rail(s, "현황", "05 / 05")
heading(s, "현재 상태와 다음 단계")

rows = [("항목", "담당", "상태"),
        ("태블릿 앱 · 서명 APK", "브릿지 팀", "완료"),
        ("브릿지 단일 실행 파일 · 배포본", "브릿지 팀", "완료"),
        ("문진 전송 · 워크리스트 조회 검증", "브릿지 팀", "완료"),
        ("연동 규격 문서 · 운영자 가이드", "브릿지 팀", "완료"),
        ("워크리스트 태그 보강 · 문진 전량 탑재", "브릿지 팀", "완료"),
        ("MPPS SCP 구현", "브릿지 팀", "완료"),
        ("매뉴얼 · 실사 160건 분석", "브릿지 팀", "완료"),
        ("MPPS 설정 · 실기 확인", "MDVizio-X", "오늘"),
        ("워크리스트 실조회 검증", "MDVizio-X", "미착수"),
        ("ai.dcm 신원 소실 · AI 임계값", "AI 담당", "신규"),
        ("실물 태블릿 · 포트 확인 · 운영자 교육", "현장", "준비 중")]
shape = table(s, M, Cm(4.0), W - 2 * M, rows, col_w=[0.56, 0.22, 0.22], size=9)
for r in range(1, len(rows)):
    cell = shape.table.cell(r, 2)
    run = cell.text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    run.font.color.rgb = BRAND if rows[r][2] == "완료" else WARN

cw3 = int((W - 2 * M - Cm(1.0)) / 3)
card(s, M, Cm(12.6), cw3, Cm(3.8), "오늘 켤 것",
     "설정 > NETWORK > MPPS에 워크리스트와 같은 주소·포트·AE 입력. 첫 환자에서 상태가 자동으로 넘어가는지 함께 확인.",
     fill=WASH)
card(s, M + cw3 + Cm(0.5), Cm(12.6), cw3, Cm(3.8), "현장 운영",
     "스터디당 72MB · 150명이면 하루 10.8GB. PACS 서버 PC를 들고 다니지만 네트워크 문제로 못 쓰는 날이 잦아, 그런 날은 전부 랩톱에 쌓임.")
card(s, M + 2 * (cw3 + Cm(0.5)), Cm(12.6), cw3, Cm(3.8), "별도 이슈",
     "실사 160건에서 AI 파이프라인 문제 3건 확인. 연동 안건이 아니라 별건이라 문서로 정리해 뒀음 — docs/field-data-findings.md",
     fill=WARNW, title_color=WARN)

text(s, M, Cm(16.8), W - 2 * M, Cm(0.6),
     "문진 수집만 놓고 보면 파일럿은 즉시 가능. 상세 규격은 docs/integration-spec.md 참조.",
     size=9.5, color=MUTED)

out = r"D:\MedMissionBridge\docs\medmission-briefing.pptx"
prs.save(out)
print("saved", out, len(prs.slides._sldIdLst), "slides")
